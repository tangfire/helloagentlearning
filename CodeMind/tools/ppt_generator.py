"""
PPT 生成器

将诊断分析报告转换为 PowerPoint 演示文稿
"""

import os
from typing import List, Dict, Any
from datetime import datetime
from pathlib import Path


class PPTGenerator:
    """
    PPT 生成器
    
    使用 python-pptx 库生成 PowerPoint 文件
    """
    
    def __init__(self, output_dir: str = None):
        """
        初始化 PPT 生成器
        
        Args:
            output_dir: 输出目录，默认为临时目录
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RgbColor
            from pptx.enum.text import PP_ALIGN
            self.Presentation = Presentation
            self.Inches = Inches
            self.Pt = Pt
            self.RgbColor = RgbColor
            self.PP_ALIGN = PP_ALIGN
            self.available = True
        except ImportError:
            print("⚠️ 未安装 python-pptx 库，PPT 生成功能不可用")
            self.available = False
        
        self.output_dir = output_dir or os.path.join(os.path.dirname(__file__), 'output', 'reports')
        os.makedirs(self.output_dir, exist_ok=True)
    
    def generate_from_report(self, report_data: Dict[str, Any], 
                            title: str = None,
                            template_path: str = None) -> str:
        """
        从报告数据生成 PPT
        
        Args:
            report_data: 报告数据，包含 analysis, conclusions, recommendations
            title: PPT 标题
            template_path: PPT 模板路径（可选）
            
        Returns:
            生成的 PPT 文件路径
        """
        if not self.available:
            raise RuntimeError("PPT 生成器不可用，请先安装 python-pptx 库")
        
        # 创建或加载模板
        if template_path and os.path.exists(template_path):
            prs = self.Presentation(template_path)
        else:
            prs = self.Presentation()
        
        # 1. 封面页
        self._add_title_slide(prs, title or "企业诊断报告")
        
        # 2. 目录页
        self._add_table_of_contents(prs)
        
        # 3. 企业概况
        if 'enterprise' in report_data:
            self._add_enterprise_overview(prs, report_data['enterprise'])
        
        # 4. 分析内容
        if 'analysis' in report_data:
            self._add_analysis_content(prs, report_data['analysis'])
        
        # 5. 主要结论
        if 'conclusions' in report_data and report_data['conclusions']:
            self._add_conclusions(prs, report_data['conclusions'])
        
        # 6. 改进建议
        if 'recommendations' in report_data and report_data['recommendations']:
            self._add_recommendations(prs, report_data['recommendations'])
        
        # 7. 结束页
        self._add_ending_slide(prs)
        
        # 保存文件
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"诊断报告_{timestamp}.pptx"
        file_path = os.path.join(self.output_dir, filename)
        
        prs.save(file_path)
        
        print(f"✅ PPT 生成成功：{file_path}")
        return file_path
    
    def _add_title_slide(self, prs, title: str):
        """添加封面页"""
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(slide_layout)
        
        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]
        
        title_placeholder.text = title
        subtitle_placeholder.text = f"生成时间：{datetime.now().strftime('%Y年%m月%d日')}"
    
    def _add_table_of_contents(self, prs):
        """添加目录页"""
        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "目录"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        sections = [
            "企业概况",
            "经营分析",
            "问题诊断",
            "主要结论",
            "改进建议"
        ]
        
        for section in sections:
            p = tf.add_paragraph()
            p.text = f"• {section}"
            p.level = 0
            p.font.size = self.Pt(24)
    
    def _add_enterprise_overview(self, prs, enterprise):
        """添加企业概况页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "企业概况"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # 企业信息
        info_items = [
            f"企业名称：{enterprise.name}",
            f"企业编码：{enterprise.code}",
            f"所属行业：{enterprise.industry or '未指定'}",
            f"企业规模：{enterprise.scale or '未指定'}"
        ]
        
        for item in info_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 0
            p.font.size = self.Pt(20)
        
        # 企业描述
        if hasattr(enterprise, 'description') and enterprise.description:
            p = tf.add_paragraph()
            p.text = f"\n企业描述："
            p.font.bold = True
            p.font.size = self.Pt(20)
            
            p = tf.add_paragraph()
            p.text = enterprise.description[:200]
            p.level = 1
            p.font.size = self.Pt(18)
    
    def _add_analysis_content(self, prs, analysis_text: str):
        """添加分析内容页"""
        # 将长文本分割成多页
        lines = analysis_text.split('\n')
        
        current_slide = None
        current_tf = None
        line_count = 0
        slide_index = 0
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 每 8 行创建新幻灯片
            if line_count % 8 == 0:
                slide_layout = prs.slide_layouts[1]
                current_slide = prs.slides.add_slide(slide_layout)
                
                title = current_slide.shapes.title
                title.text = f"经营分析 ({slide_index + 1})"
                
                current_tf = current_slide.placeholders[1].text_frame
                current_tf.clear()
                slide_index += 1
            
            p = current_tf.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = self.Pt(18)
            
            line_count += 1
    
    def _add_conclusions(self, prs, conclusions: List[str]):
        """添加结论页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "主要结论"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for conclusion in conclusions:
            p = tf.add_paragraph()
            p.text = f"✓ {conclusion}"
            p.level = 0
            p.font.size = self.Pt(20)
            p.space_after = self.Pt(12)
    
    def _add_recommendations(self, prs, recommendations: List[str]):
        """添加建议页"""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "改进建议"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        for i, rec in enumerate(recommendations, 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {rec}"
            p.level = 0
            p.font.size = self.Pt(20)
            p.space_after = self.Pt(12)
    
    def _add_ending_slide(self, prs):
        """添加结束页"""
        slide_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加居中的感谢文字
        left = self.Inches(2)
        top = self.Inches(3)
        width = self.Inches(6)
        height = self.Inches(2)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = "感谢聆听"
        p.alignment = self.PP_ALIGN.CENTER
        p.font.size = self.Pt(48)
        p.font.bold = True


# 便捷函数
def generate_diagnostic_ppt(report_data: Dict[str, Any], 
                           title: str = None,
                           output_dir: str = None) -> str:
    """
    便捷函数：生成诊断报告 PPT
    
    Args:
        report_data: 报告数据
        title: PPT 标题
        output_dir: 输出目录
        
    Returns:
        PPT 文件路径
    """
    generator = PPTGenerator(output_dir=output_dir)
    return generator.generate_from_report(report_data, title=title)
